import tempfile
from pathlib import Path

import fitz
from pptx import Presentation

from speechpt.coherence.document_parser import parse_pdf, parse_ppt


def test_parse_pdf_extracts_title_and_text():
    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((72, 72), "Title\n- point1\n- point2")
        doc.save(tmp.name)
        doc.close()
        slides = parse_pdf(Path(tmp.name))
    assert slides[0].title == "Title"
    assert "point1" in slides[0].text
    assert isinstance(slides[0].visual_items, list)
    assert isinstance(slides[0].visual_captions, list)
    Path(tmp.name).unlink()


def test_parse_ppt_reads_shapes():
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "My Title"
    slide.placeholders[1].text = "- bullet"
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        prs.save(tmp.name)
        slides = parse_ppt(Path(tmp.name))
    assert slides[0].title == "My Title"
    assert "bullet" in slides[0].bullet_points[0]
    assert isinstance(slides[0].visual_items, list)
    assert isinstance(slides[0].visual_captions, list)
    Path(tmp.name).unlink()


def test_parse_pdf_uses_first_non_bullet_line_as_title():
    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((72, 72), "- intro\n1. point1\nReal Title\n2. point2")
        doc.save(tmp.name)
        doc.close()
        slides = parse_pdf(Path(tmp.name))
    assert slides[0].title == "Real Title"
    assert "point1" in slides[0].bullet_points
    assert "point2" in slides[0].bullet_points
    Path(tmp.name).unlink()


def test_parse_ppt_falls_back_to_first_non_bullet_text_for_title():
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    textbox = slide.shapes.add_textbox(0, 0, 3000000, 800000)
    textbox.text_frame.text = "Fallback Title"
    body = slide.shapes.add_textbox(0, 900000, 3000000, 1200000)
    body.text_frame.text = "1. bullet one"
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        prs.save(tmp.name)
        slides = parse_ppt(Path(tmp.name))
    assert slides[0].title == "Fallback Title"
    assert "bullet one" in slides[0].bullet_points
    Path(tmp.name).unlink()


def test_parse_pdf_skips_icon_only_title_and_reads_icon_prefixed_bullets():
    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((72, 72), "\uf274\nOUR MISSION\n\uf00cFirst point\n\uf00cSecond point")
        doc.save(tmp.name)
        doc.close()
        slides = parse_pdf(Path(tmp.name))
    assert slides[0].title == "OUR MISSION"
    assert "First point" in slides[0].bullet_points
    assert "Second point" in slides[0].bullet_points
    assert "\uf00c" not in slides[0].text
    Path(tmp.name).unlink()
