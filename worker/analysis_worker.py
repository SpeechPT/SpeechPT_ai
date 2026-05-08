"""SpeechPT 분석 Worker (Step 9: 실제 SageMaker 통합).

Step 8까지: mock 결과 반환
Step 9~  : 실제 SageMaker 통합 EP 호출 (STT → CE → AE)

흐름:
1. SQS 폴링
2. analyses UPDATE → status=running
3. document parser (Worker CPU) → slide keypoints
4. SageMaker invoke STT → words
5. SageMaker invoke CE → coverage per slide
6. SageMaker invoke AE → segment scores
7. report_json 합산 → analysis_results INSERT
8. analyses UPDATE → status=done
"""
from __future__ import annotations

import json
import os
import time
import traceback
import uuid
from datetime import datetime, timezone
from typing import Any, Dict, List

import boto3
from sqlalchemy import create_engine, text
from sqlalchemy.orm import sessionmaker

from worker.sagemaker_client import invoke as sm_invoke
from speechpt.coherence.keypoint_extractor import Keypoint
from speechpt.coherence.auto_aligner import auto_align_slides
from speechpt.coherence.transcript_aligner import align_transcript

# ──────────────────────────────────────────────────────────────
# 환경 변수
# ──────────────────────────────────────────────────────────────
REGION = os.environ.get("AWS_REGION", "ap-northeast-2")
QUEUE_URL = os.environ["ANALYSIS_QUEUE_URL"]
RESULTS_BUCKET = os.environ["S3_BUCKET_RESULTS"]
DATABASE_URL = os.environ["DATABASE_URL"]
USE_MOCK = os.environ.get("WORKER_USE_MOCK", "false").lower() == "true"

WORKER_ID = f"worker-{os.environ.get('HOSTNAME', uuid.uuid4().hex[:8])}"

# ──────────────────────────────────────────────────────────────
# AWS / DB 클라이언트
# ──────────────────────────────────────────────────────────────
sqs = boto3.client("sqs", region_name=REGION)
s3 = boto3.client("s3", region_name=REGION)

engine = create_engine(DATABASE_URL, pool_pre_ping=True, pool_recycle=3600)
SessionLocal = sessionmaker(bind=engine)


def log(msg: str) -> None:
    print(f"[{datetime.utcnow().isoformat()}][{WORKER_ID}] {msg}", flush=True)


# ──────────────────────────────────────────────────────────────
# DB 헬퍼
# ──────────────────────────────────────────────────────────────
def update_analysis(analysis_id: str, **fields: Any) -> None:
    if not fields:
        return
    cols = ", ".join(f"{k} = :{k}" for k in fields)
    with SessionLocal() as db:
        db.execute(
            text(f"UPDATE analyses SET {cols} WHERE analysis_id = :id"),
            {**fields, "id": analysis_id},
        )
        db.commit()


def upsert_analysis_result(
    analysis_id: str,
    *,
    content_coverage: int,
    delivery_stability: int,
    pacing_score: int,
    overall_score: int,
    severity_json: dict,
    report_json: dict,
    result_blob_s3_key: str,
) -> None:
    sql = text("""
        INSERT INTO analysis_results (
            analysis_id, content_coverage, delivery_stability, pacing_score,
            overall_score, severity_json, report_json, result_blob_s3_key
        ) VALUES (
            :analysis_id, :content_coverage, :delivery_stability, :pacing_score,
            :overall_score, CAST(:severity_json AS JSONB), CAST(:report_json AS JSONB),
            :result_blob_s3_key
        )
        ON CONFLICT (analysis_id) DO UPDATE SET
            content_coverage = EXCLUDED.content_coverage,
            delivery_stability = EXCLUDED.delivery_stability,
            pacing_score = EXCLUDED.pacing_score,
            overall_score = EXCLUDED.overall_score,
            severity_json = EXCLUDED.severity_json,
            report_json = EXCLUDED.report_json,
            result_blob_s3_key = EXCLUDED.result_blob_s3_key
    """)
    with SessionLocal() as db:
        db.execute(sql, {
            "analysis_id": analysis_id,
            "content_coverage": content_coverage,
            "delivery_stability": delivery_stability,
            "pacing_score": pacing_score,
            "overall_score": overall_score,
            "severity_json": json.dumps(severity_json, ensure_ascii=False),
            "report_json": json.dumps(report_json, ensure_ascii=False),
            "result_blob_s3_key": result_blob_s3_key,
        })
        db.commit()


# ──────────────────────────────────────────────────────────────
# Document parsing (Worker CPU에서 수행)
# ──────────────────────────────────────────────────────────────
def parse_document_to_keypoints(doc_uri: str) -> List[Dict[str, Any]]:
    """슬라이드 → 키포인트 리스트.

    임시 구현: PDF/PPTX 다운로드 후 텍스트 추출. 슬라이드별 1개 키포인트.
    Step 10에서 keypoint_extractor로 정교화.
    """
    import tempfile

    # 다운로드
    bucket, key = doc_uri.replace("s3://", "").split("/", 1)
    suffix = os.path.splitext(key)[1].lower()
    fd, local = tempfile.mkstemp(suffix=suffix)
    os.close(fd)
    s3.download_file(bucket, key, local)

    slides: List[Dict[str, Any]] = []

    if suffix == ".pdf":
        try:
            import fitz  # PyMuPDF — Worker requirements에 추가 필요
            doc = fitz.open(local)
            for i, page in enumerate(doc):
                text = page.get_text().strip()
                first_line = text.split("\n", 1)[0] if text else f"Slide {i+1}"
                slides.append({
                    "slide_id": i + 1,
                    "title": first_line[:100],
                    "keypoints": [first_line] if first_line else [f"Slide {i+1}"],
                    "full_text": text,
                })
        except ImportError:
            log("PyMuPDF not installed — fallback to single keypoint")
            slides.append({"slide_id": 1, "title": "Document",
                          "keypoints": ["Document content"], "full_text": ""})
    elif suffix in (".ppt", ".pptx"):
        try:
            from pptx import Presentation
            prs = Presentation(local)
            for i, slide in enumerate(prs.slides):
                text_parts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                text_parts.append(t)
                title = text_parts[0] if text_parts else f"Slide {i+1}"
                slides.append({
                    "slide_id": i + 1,
                    "title": title[:100],
                    "keypoints": text_parts[:5] or [title],
                    "full_text": "\n".join(text_parts),
                })
        except ImportError:
            log("python-pptx not installed — fallback to single keypoint")
            slides.append({"slide_id": 1, "title": "Document",
                          "keypoints": ["Document content"], "full_text": ""})
    else:
        slides.append({"slide_id": 1, "title": "Unknown format",
                      "keypoints": ["Document content"], "full_text": ""})

    os.unlink(local)
    return slides


def chunk_transcript(words: List[Dict[str, Any]], duration_sec: float = 10.0) -> List[Dict[str, Any]]:
    """단어 list를 시간 청크로 묶음 (CE 입력용)."""
    if not words:
        return []
    chunks = []
    current_start = float(words[0].get("start", 0.0))
    current_words = []
    for w in words:
        w_start = float(w.get("start", 0.0))
        if current_words and (w_start - current_start) >= duration_sec:
            chunks.append({
                "start": current_start,
                "end": float(current_words[-1].get("end", current_start)),
                "text": " ".join(str(cw.get("word", "")) for cw in current_words),
            })
            current_words = [w]
            current_start = w_start
        else:
            current_words.append(w)
    if current_words:
        chunks.append({
            "start": current_start,
            "end": float(current_words[-1].get("end", current_start)),
            "text": " ".join(str(cw.get("word", "")) for cw in current_words),
        })
    return chunks


# ──────────────────────────────────────────────────────────────
# 실제 분석 — SageMaker 호출 사슬
# ──────────────────────────────────────────────────────────────
def run_real_analysis(payload: dict) -> dict:
    """STT → CE → AE 파이프라인. SageMaker 통합 EP 사용."""
    audio_uri = f"s3://{payload['audio']['bucket']}/{payload['audio']['object_key']}"
    doc_uri = f"s3://{payload['document']['bucket']}/{payload['document']['object_key']}"

    # 1. 문서 파싱 (Worker CPU)
    log("  parsing document...")
    slides = parse_document_to_keypoints(doc_uri)
    log(f"  → {len(slides)} slides")

    # 2. STT (SageMaker)
    log("  invoking STT...")
    stt_result = sm_invoke("stt", {"audio_s3": audio_uri}, timeout_sec=900)
    words = stt_result.get("words", [])
    transcript = stt_result.get("text", "")
    log(f"  → {len(words)} words")

    # 3. CE (SageMaker) — 키포인트 vs 트랜스크립트 청크
    log("  invoking CE...")
    chunks = chunk_transcript(words)
    all_keypoints: List[str] = []
    for s in slides:
        all_keypoints.extend(s.get("keypoints", []))
    if chunks and all_keypoints:
        ce_result = sm_invoke("ce", {
            "keypoints": all_keypoints,
            "transcript_chunks": [c["text"] for c in chunks],
            "threshold": 0.55,
        }, timeout_sec=300)
    else:
        ce_result = {"coverage": 0.0, "covered_mask": []}

    # 4. AE (SageMaker) — 슬라이드별 segment
    log("  invoking auto_aligner for accurate segments...")
    slide_keypoints = []
    for s in slides:
        kps = []
        if s.get("title"):
            kps.append(Keypoint(text=s["title"], importance=1.0, source="title"))
        for text_kp in s.get("keypoints", []):
            if text_kp != s.get("title"):
                kps.append(Keypoint(text=text_kp, importance=0.8, source="bullet"))
        slide_keypoints.append(kps)
        
    try:
        align_result = auto_align_slides(
            slide_keypoints=slide_keypoints,
            words=words,
            model_name="jhgan/ko-sroberta-multitask"
        )
        boundaries = align_result.final_boundaries
        transcript_segments = align_transcript(words, boundaries)
        
        segments = []
        for i, (s, tseg) in enumerate(zip(slides, transcript_segments)):
            segments.append({
                "slide_id": s["slide_id"],
                "start_sec": tseg.start_sec,
                "end_sec": tseg.end_sec,
            })
    except Exception as e:
        log(f"  auto_aligner failed: {e}. Falling back to uniform split.")
        total_dur = max((w.get("end", 0) for w in words), default=0)
        if total_dur <= 0 and slides:
            total_dur = len(slides) * 30
        seg_dur = total_dur / max(len(slides), 1)
        segments = [
            {
                "slide_id": s["slide_id"],
                "start_sec": i * seg_dur,
                "end_sec": (i + 1) * seg_dur,
            }
            for i, s in enumerate(slides)
        ]

    log("  invoking AE...")
    ae_result = sm_invoke("ae", {
        "audio_s3": audio_uri,
        "segments": segments,
        "chunk_sec": 30,
    }, timeout_sec=900)

    # 5. 점수 합산
    coverage_pct = int(round(ce_result.get("coverage", 0.0) * 100))
    ae_segments = ae_result.get("segments", [])

    # AE 파생 점수
    valid_segs = [s for s in ae_segments if s.get("scores")]
    if valid_segs:
        avg_overall = sum(s["scores"]["overall_delivery"] for s in valid_segs) / len(valid_segs)
        delivery_stability = max(0, min(100, int(round(avg_overall * 100))))
        rates = [s["scores"]["speech_rate"] for s in valid_segs]
        avg_rate = sum(rates) / len(rates)
        # 이상적인 말 속도 ~2.5 음절/초 가정
        rate_diff = abs(avg_rate - 2.5)
        pacing_score = max(0, min(100, int(round(100 - rate_diff * 30))))
    else:
        delivery_stability = 50
        pacing_score = 50

    overall_score = int(round((coverage_pct + delivery_stability + pacing_score) / 3))

    # 6. report_json 구성 (BE 응답 스키마와 일치)
    sections = []
    for i, (s, ae_seg) in enumerate(zip(slides, ae_segments)):
        ae_score = 50
        feedback = "분석 데이터 부족"
        if ae_seg.get("scores"):
            sc = ae_seg["scores"]
            ae_score = max(0, min(100, int(round(sc.get("overall_delivery", 0.5) * 100))))
            issues = []
            if sc.get("energy_drop_prob", 0) > 0.5:
                issues.append("에너지 저하")
            if sc.get("pitch_shift_prob", 0) > 0.5:
                issues.append("피치 흔들림")
            sr = sc.get("speech_rate", 0)
            if sr > 3.5:
                issues.append("말 속도 빠름")
            elif sr < 1.5:
                issues.append("말 속도 느림")
            feedback = ", ".join(issues) if issues else "안정적인 발화"

        sections.append({
            "section_index": i + 1,
            "title": s.get("title", f"슬라이드 {i+1}"),
            "start_time_sec": int(ae_seg.get("start_sec", 0)),
            "end_time_sec": int(ae_seg.get("end_sec", 0)),
            "score": ae_score,
            "feedback": feedback,
        })

    strengths: List[Dict[str, str]] = []
    improvements: List[Dict[str, str]] = []
    if coverage_pct >= 70:
        strengths.append({"text": f"슬라이드 키포인트 커버리지가 {coverage_pct}%로 양호합니다."})
    else:
        improvements.append({"text": f"슬라이드 키포인트 중 {100-coverage_pct}%가 발화에서 충분히 다뤄지지 않았습니다."})
    if delivery_stability >= 70:
        strengths.append({"text": "전반적인 발화가 안정적입니다."})
    else:
        improvements.append({"text": "발화 안정성을 높이기 위해 호흡 조절이 필요합니다."})
    if pacing_score < 60:
        improvements.append({"text": "말 속도가 일정하지 않아 청중 이해도가 떨어질 수 있습니다."})

    report = {
        "summary": (f"전체 점수 {overall_score}점 — 내용 일치도 {coverage_pct}, "
                    f"발화 안정성 {delivery_stability}, 페이싱 {pacing_score}."),
        "strengths": strengths,
        "improvements": improvements,
        "sections": sections,
        "transcript": transcript[:5000],
    }

    return {
        "scores": {
            "content_coverage": coverage_pct,
            "delivery_stability": delivery_stability,
            "pacing_score": pacing_score,
            "overall_score": overall_score,
        },
        "severity": {
            "high_severity_count": sum(1 for s in sections if s["score"] < 50),
            "medium_severity_count": sum(1 for s in sections if 50 <= s["score"] < 70),
            "low_severity_count": sum(1 for s in sections if s["score"] >= 70),
        },
        "report": report,
    }


# ──────────────────────────────────────────────────────────────
# Mock 분석 (USE_MOCK=true 일 때만)
# ──────────────────────────────────────────────────────────────
def run_mock_analysis(payload: dict) -> dict:
    return {
        "scores": {
            "content_coverage": 75, "delivery_stability": 80,
            "pacing_score": 72, "overall_score": 76,
        },
        "severity": {"high_severity_count": 0, "medium_severity_count": 1, "low_severity_count": 2},
        "report": {
            "summary": "MOCK 결과입니다. WORKER_USE_MOCK=false로 실제 SageMaker 호출 활성화.",
            "strengths": [{"text": "MOCK: 안정적 발화"}],
            "improvements": [{"text": "MOCK: 후반부 속도"}],
            "sections": [],
        },
    }


# ──────────────────────────────────────────────────────────────
# 메인 처리
# ──────────────────────────────────────────────────────────────
def process(payload: dict) -> None:
    aid = payload["analysis_id"]
    log(f"START analysis {aid} (USE_MOCK={USE_MOCK})")

    update_analysis(
        aid,
        status="running", progress=10, stage="ingest",
        worker_id=WORKER_ID,
        started_at=datetime.now(timezone.utc),
    )

    try:
        if USE_MOCK:
            time.sleep(3)
            update_analysis(aid, progress=40, stage="analyzing")
            time.sleep(3)
            update_analysis(aid, progress=70)
            time.sleep(3)
            result = run_mock_analysis(payload)
        else:
            update_analysis(aid, progress=20, stage="analyzing")
            result = run_real_analysis(payload)
            update_analysis(aid, progress=90)
    except Exception:
        log(f"ANALYSIS ERROR for {aid}:\n{traceback.format_exc()}")
        raise

    # S3에 백업
    result_blob_key = f"results/{aid}.json"
    s3.put_object(
        Bucket=RESULTS_BUCKET,
        Key=result_blob_key,
        Body=json.dumps(result, ensure_ascii=False).encode("utf-8"),
        ContentType="application/json",
    )

    upsert_analysis_result(
        aid,
        content_coverage=result["scores"]["content_coverage"],
        delivery_stability=result["scores"]["delivery_stability"],
        pacing_score=result["scores"]["pacing_score"],
        overall_score=result["scores"]["overall_score"],
        severity_json=result["severity"],
        report_json=result["report"],
        result_blob_s3_key=result_blob_key,
    )

    update_analysis(
        aid, status="done", progress=100, stage="finished",
        finished_at=datetime.now(timezone.utc),
    )
    log(f"DONE analysis {aid}")


def main() -> None:
    log(f"Worker started")
    log(f"  QUEUE_URL = {QUEUE_URL}")
    log(f"  RESULTS_BUCKET = {RESULTS_BUCKET}")
    log(f"  USE_MOCK = {USE_MOCK}")
    log(f"  AE_ENDPOINT_NAME = {os.environ.get('AE_ENDPOINT_NAME', 'speechpt-unified-async')}")

    while True:
        try:
            resp = sqs.receive_message(
                QueueUrl=QUEUE_URL,
                MaxNumberOfMessages=1,
                WaitTimeSeconds=20,
                VisibilityTimeout=1800,    # SageMaker 호출 ~3분 + 여유
            )
            for msg in resp.get("Messages", []):
                try:
                    payload = json.loads(msg["Body"])
                except json.JSONDecodeError as exc:
                    log(f"INVALID JSON: {exc}")
                    sqs.delete_message(QueueUrl=QUEUE_URL, ReceiptHandle=msg["ReceiptHandle"])
                    continue

                try:
                    process(payload)
                except Exception as exc:
                    log(f"FAILED {payload.get('analysis_id')}: {type(exc).__name__}: {exc}")
                    aid = payload.get("analysis_id")
                    if aid:
                        try:
                            update_analysis(
                                aid, status="failed",
                                error_code=type(exc).__name__,
                                error_message=str(exc)[:500],
                                finished_at=datetime.now(timezone.utc),
                            )
                        except Exception as db_exc:
                            log(f"db update failed: {db_exc}")
                finally:
                    sqs.delete_message(QueueUrl=QUEUE_URL, ReceiptHandle=msg["ReceiptHandle"])
        except KeyboardInterrupt:
            log("Worker shutting down")
            break
        except Exception as exc:
            log(f"polling error: {exc}, sleep 5s")
            time.sleep(5)


if __name__ == "__main__":
    main()
