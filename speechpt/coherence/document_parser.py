"""문서(PDF/PPT)에서 슬라이드별 텍스트와 시각 요소를 구조화하여 추출하는 모듈."""
from __future__ import annotations

import argparse
import logging
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)
logging.basicConfig(level=logging.INFO)

_BULLET_PREFIX_RE = re.compile(r"^(?:[-•·*●▪◦‣]|[0-9]+[.)]|[A-Za-z][.)])\s*")


@dataclass
class VisualItem:
    """슬라이드 내 시각 요소 메타데이터."""

    item_id: str
    slide_id: int
    item_type: str
    source: str  # "pdf" | "ppt"
    bbox: List[float] | None = None
    raw_text: str = ""
    confidence: float = 0.0
    metadata: Dict[str, Any] = field(default_factory=dict)


@dataclass
class SlideContent:
    """슬라이드 단위 텍스트/시각 요소 컨테이너."""

    slide_id: int
    text: str
    title: str
    bullet_points: List[str]
    visual_captions: List[str] = field(default_factory=list)
    visual_items: List[VisualItem] = field(default_factory=list)


@dataclass
class PdfLine:
    text: str
    font_size: float
    y0: float
    x0: float


def _strip_bullet_prefix(text: str) -> str:
    stripped = text.strip()
    if stripped and 0xF000 <= ord(stripped[0]) <= 0xF8FF:
        return stripped[1:].strip()
    return _BULLET_PREFIX_RE.sub("", stripped).strip()


def _normalize_pdf_text(text: str) -> str:
    stripped = " ".join(text.strip().split())
    if not stripped:
        return stripped
    tokens = stripped.split(" ")
    if len(tokens) >= 4:
        short_tokens = sum(1 for token in tokens if len(token) <= 2)
        if short_tokens / max(1, len(tokens)) >= 0.75:
            return "".join(tokens)
    return stripped


def _is_iconish_line(text: str) -> bool:
    stripped = text.strip()
    if not stripped:
        return False
    return all(0xF000 <= ord(ch) <= 0xF8FF for ch in stripped)


def _is_metricish_line(text: str) -> bool:
    stripped = text.strip()
    if not stripped:
        return False
    compact = stripped.replace(" ", "")
    metric_tokens = {"%", "H", "년차", "건", "초", "명", "억+", "AI"}
    if any(token in compact for token in metric_tokens) and len(compact) <= 6:
        return True
    return compact.replace(".", "").replace("+", "").isdigit()


def _title_score(text: str) -> float:
    stripped = text.strip()
    if not stripped:
        return -1.0
    score = float(len(stripped))
    if _is_iconish_line(stripped):
        return -1.0
    if _is_metricish_line(stripped):
        score -= 20.0
    if _looks_like_bullet(stripped):
        score -= 25.0
    if any(ch.isalpha() for ch in stripped):
        score += 2.0
    if any("\uac00" <= ch <= "\ud7a3" for ch in stripped):
        score += 2.0
    if stripped.isupper():
        score += 1.0
    return score


def _pdf_title_score(line: PdfLine) -> float:
    stripped = line.text.strip()
    score = _title_score(stripped)
    if score < 0:
        return score
    score += min(line.font_size, 40.0) * 2.2
    score += max(0.0, 180.0 - min(line.y0, 180.0)) * 0.12
    if len(stripped) > 60:
        score -= 12.0
    if len(stripped) < 3:
        score -= 6.0
    if line.font_size >= 24.0 and line.y0 <= 110.0:
        score += 18.0
    elif line.font_size >= 18.0 and line.y0 <= 140.0:
        score += 8.0
    return score


def _looks_like_bullet(text: str) -> bool:
    stripped = text.strip()
    if not stripped:
        return False
    if len(stripped) > 1 and 0xF000 <= ord(stripped[0]) <= 0xF8FF:
        return True
    if _BULLET_PREFIX_RE.match(stripped):
        return True
    # PPT paragraphs often lose explicit markers but remain short noun-phrase lines.
    return False


def _extract_bullets(lines: List[str]) -> List[str]:
    bullets = []
    for line in lines:
        stripped = line.strip()
        if not stripped:
            continue
        if _looks_like_bullet(stripped):
            normalized = _strip_bullet_prefix(stripped)
            if normalized:
                bullets.append(normalized)
    return bullets


def _normalize_content_line(text: str) -> str:
    stripped = _normalize_pdf_text(text)
    if not stripped or _is_iconish_line(stripped):
        return ""
    if _looks_like_bullet(stripped):
        return _strip_bullet_prefix(stripped)
    if stripped and 0xF000 <= ord(stripped[0]) <= 0xF8FF:
        return stripped[1:].strip()
    return stripped


def _select_title(lines: List[str]) -> str:
    candidates = []
    for idx, line in enumerate(lines[:8]):
        stripped = line.strip()
        if not stripped or _is_iconish_line(stripped):
            continue
        candidates.append((idx, _title_score(stripped), stripped))
    if not candidates:
        return lines[0].strip() if lines else ""
    candidates.sort(key=lambda item: (item[1], -item[0]), reverse=True)
    return candidates[0][2]


def _extract_pdf_lines(page: fitz.Page) -> List[PdfLine]:
    text_dict = page.get_text("dict")
    extracted: List[PdfLine] = []
    for block in text_dict.get("blocks", []):
        if block.get("type") != 0:
            continue
        for line in block.get("lines", []):
            spans = line.get("spans", [])
            parts = []
            sizes = []
            y_positions = []
            x_positions = []
            for span in spans:
                text = str(span.get("text", "")).strip()
                if not text:
                    continue
                parts.append(text)
                sizes.append(float(span.get("size", 0.0)))
                bbox = span.get("bbox", [0.0, 0.0, 0.0, 0.0])
                y_positions.append(float(bbox[1]))
                x_positions.append(float(bbox[0]))
            if not parts:
                continue
            extracted.append(
                PdfLine(
                    text=_normalize_pdf_text(" ".join(parts)),
                    font_size=max(sizes) if sizes else 0.0,
                    y0=min(y_positions) if y_positions else 0.0,
                    x0=min(x_positions) if x_positions else 0.0,
                )
            )
    extracted.sort(key=lambda line: (line.y0, line.x0))
    return extracted


def _select_pdf_title(pdf_lines: List[PdfLine], fallback_lines: List[str]) -> str:
    candidates = []
    for idx, line in enumerate(pdf_lines):
        stripped = line.text.strip()
        if not stripped or _is_iconish_line(stripped):
            continue
        candidates.append((idx, _pdf_title_score(line), stripped))
    if not candidates:
        return _select_title(fallback_lines)
    priority = [item for item in candidates if item[1] >= 60.0]
    if priority:
        priority.sort(key=lambda item: (item[1], -item[0]), reverse=True)
        return priority[0][2]
    candidates.sort(key=lambda item: (item[1], -item[0]), reverse=True)
    return candidates[0][2]


def _summarize_visual_items(items: List[VisualItem]) -> List[str]:
    counts: Dict[str, int] = {}
    for item in items:
        counts[item.item_type] = counts.get(item.item_type, 0) + 1
    captions = []
    for item_type, count in sorted(counts.items()):
        captions.append(f"{item_type} x{count}")
    return captions


def _extract_pdf_visual_items(page: fitz.Page, slide_id: int) -> List[VisualItem]:
    items: List[VisualItem] = []

    for idx, image_info in enumerate(page.get_images(full=True), start=1):
        xref = image_info[0]
        rects = page.get_image_rects(xref)
        bbox = None
        if rects:
            r = rects[0]
            bbox = [float(r.x0), float(r.y0), float(r.x1), float(r.y1)]
        items.append(
            VisualItem(
                item_id=f"s{slide_id}_img_{idx}",
                slide_id=slide_id,
                item_type="image",
                source="pdf",
                bbox=bbox,
                metadata={"xref": xref},
            )
        )

    drawings = page.get_drawings()
    if drawings:
        items.append(
            VisualItem(
                item_id=f"s{slide_id}_chart_candidate_1",
                slide_id=slide_id,
                item_type="chart_candidate",
                source="pdf",
                metadata={"drawings_count": len(drawings)},
            )
        )

    try:
        tables = page.find_tables()
        if getattr(tables, "tables", None):
            for idx, _ in enumerate(tables.tables, start=1):
                items.append(
                    VisualItem(
                        item_id=f"s{slide_id}_table_{idx}",
                        slide_id=slide_id,
                        item_type="table",
                        source="pdf",
                    )
                )
    except Exception:
        # PyMuPDF 버전에 따라 find_tables 미지원일 수 있어 무시한다.
        pass

    return items


def parse_pdf(path: Path) -> List[SlideContent]:
    doc = fitz.open(path)
    slides: List[SlideContent] = []
    try:
        for i, page in enumerate(doc, start=1):
            pdf_lines = _extract_pdf_lines(page)
            raw_lines = [line.text for line in pdf_lines if line.text.strip()]
            lines = []
            for line in pdf_lines:
                normalized = _normalize_content_line(line.text)
                if normalized:
                    lines.append(normalized)
            title = _select_pdf_title(pdf_lines, lines)
            bullet_points = _extract_bullets(raw_lines)
            visual_items = _extract_pdf_visual_items(page, i)
            slides.append(
                SlideContent(
                    slide_id=i,
                    text="\n".join(lines),
                    title=title,
                    bullet_points=bullet_points,
                    visual_captions=_summarize_visual_items(visual_items),
                    visual_items=visual_items,
                )
            )
    finally:
        doc.close()
    return slides


def _shape_to_item_type(shape) -> str | None:
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return "image"
    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        return "table"
    if shape.shape_type == MSO_SHAPE_TYPE.CHART:
        return "chart"
    smart_art = getattr(MSO_SHAPE_TYPE, "SMART_ART", None)
    if smart_art is not None and shape.shape_type == smart_art:
        return "diagram"
    return None


def _extract_ppt_visual_items(slide, slide_id: int) -> List[VisualItem]:
    items: List[VisualItem] = []
    for idx, shape in enumerate(slide.shapes, start=1):
        item_type = _shape_to_item_type(shape)
        if item_type is None:
            continue

        bbox = [float(shape.left), float(shape.top), float(shape.left + shape.width), float(shape.top + shape.height)]
        metadata: Dict[str, Any] = {"shape_id": shape.shape_id}
        if item_type == "chart" and hasattr(shape, "chart"):
            metadata["chart_type"] = str(shape.chart.chart_type)
        items.append(
            VisualItem(
                item_id=f"s{slide_id}_{item_type}_{idx}",
                slide_id=slide_id,
                item_type=item_type,
                source="ppt",
                bbox=bbox,
                metadata=metadata,
            )
        )
    return items


def parse_ppt(path: Path) -> List[SlideContent]:
    prs = Presentation(path)
    slides: List[SlideContent] = []
    for i, slide in enumerate(prs.slides, start=1):
        texts: List[str] = []
        title_text = ""
        bullet_points: List[str] = []
        visual_items = _extract_ppt_visual_items(slide, i)
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            if shape.is_placeholder and shape.placeholder_format.type in {1, 3}:
                title_text = shape.text.strip()
            for paragraph in shape.text_frame.paragraphs:
                p_text = "".join(run.text for run in paragraph.runs).strip()
                if not p_text:
                    continue
                texts.append(p_text)
                if paragraph.level > 0 or _looks_like_bullet(p_text):
                    normalized = _strip_bullet_prefix(p_text)
                    if normalized:
                        bullet_points.append(normalized)
        if not title_text:
            title_text = _select_title(texts)
        slides.append(
            SlideContent(
                slide_id=i,
                text="\n".join(texts),
                title=title_text,
                bullet_points=bullet_points,
                visual_captions=_summarize_visual_items(visual_items),
                visual_items=visual_items,
            )
        )
    return slides


def parse_document(path: str | Path) -> List[SlideContent]:
    path = Path(path)
    if not path.exists():
        raise FileNotFoundError(path)
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return parse_pdf(path)
    if suffix in {".ppt", ".pptx"}:
        return parse_ppt(path)
    raise ValueError(f"Unsupported file type: {suffix}")


def main():
    parser = argparse.ArgumentParser(description="Parse PDF/PPT into slide texts and visual item metadata")
    parser.add_argument("input", type=str, help="Path to PDF or PPTX file")
    args = parser.parse_args()

    slides = parse_document(args.input)
    for slide in slides:
        logger.info(
            "Slide %d | title=%s | bullets=%d | visuals=%d | visual_captions=%s",
            slide.slide_id,
            slide.title[:30],
            len(slide.bullet_points),
            len(slide.visual_items),
            slide.visual_captions,
        )


if __name__ == "__main__":
    main()
